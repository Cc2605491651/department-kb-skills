from __future__ import annotations

import importlib.util
from pathlib import Path
import sys
import tempfile
import unittest


REQUIRED = ("charset_normalizer", "openpyxl", "pdfplumber", "docx", "pptx", "reportlab")
HAS_DOCUMENT_RUNTIME = all(importlib.util.find_spec(name) is not None for name in REQUIRED)
SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))


@unittest.skipUnless(HAS_DOCUMENT_RUNTIME, "bundled document runtime is not active")
class DocumentFormatSmokeTests(unittest.TestCase):
    def test_primary_extractors(self) -> None:
        from docx import Document
        import openpyxl
        from pptx import Presentation
        from reportlab.pdfgen.canvas import Canvas
        from extract_all import parse_local_file

        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)

            text_path = root / "sample.txt"
            text_path.write_text("plain-text-marker", encoding="utf-8")

            docx_path = root / "sample.docx"
            document = Document()
            document.add_paragraph("docx-marker")
            document.save(docx_path)

            pptx_path = root / "sample.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            box = slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000)
            box.text = "pptx-marker with enough searchable slide text"
            presentation.save(pptx_path)

            xlsx_path = root / "sample.xlsx"
            workbook = openpyxl.Workbook()
            workbook.active["A1"] = "xlsx-marker"
            workbook.save(xlsx_path)

            pdf_path = root / "sample.pdf"
            canvas = Canvas(str(pdf_path))
            canvas.drawString(72, 720, "pdf-marker searchable business text " * 4)
            canvas.save()

            cases = [
                (text_path, "txt", "plain-text-marker"),
                (docx_path, "docx", "docx-marker"),
                (pptx_path, "pptx", "pptx-marker"),
                (xlsx_path, "xlsx", "xlsx-marker"),
                (pdf_path, "pdf", "pdf-marker"),
            ]
            for path, extension, marker in cases:
                extracted, _, _, _ = parse_local_file(path, extension, root / f"work-{extension}")
                self.assertIn(marker, extracted, extension)


if __name__ == "__main__":
    unittest.main()
